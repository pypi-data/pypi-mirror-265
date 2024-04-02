# SlideDetect - A tool for detecting slide changes in videos and saving them as PowerPoint slides.
# Copyright (C) 2024 Nindo Punturi
#
# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU General Public License for more details.
#
# You should have received a copy of the GNU General Public License
# along with this program.  If not, see <https://www.gnu.org/licenses/gpl-3.0.html>.

from pptx import Presentation
from pptx.util import Inches
import io
import cv2
import numpy as np
import argparse
from PIL import Image
from tqdm import tqdm
import os
import sys


parser = argparse.ArgumentParser(description='Slide change detection in videos.')
parser.add_argument('video_path', type=str, help='Path to the input video')
parser.add_argument('output_folder', type=str, help='Path to save extracted slides')
parser.add_argument('--threshold', type=float, default=1, help='Threshold for detecting slide changes')
args = parser.parse_args()
prs = Presentation()

def print_gpl_notice():
    print(f"""
    SlideDetect Copyright (C) {2024} Nindo Punturi
    This program comes with ABSOLUTELY NO WARRANTY; for details type `show w`.
    This is free software, and you are welcome to redistribute it
    under certain conditions; type `show c' for details.
    """)

def show_warranty():
    print(f"""
          This program is distributed in the hope that it will be useful, but WITHOUT ANY WARRANTY;
          without even the implied warranty of MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.
          See the GNU General Public License for more details <https://www.gnu.org/licenses/gpl-3.0.html>.
          """)

def show_conditions():
    print(f"""
          You are welcome to redistribute this software under certain conditions;
          see the GNU General Public License for more details at <https://www.gnu.org/licenses/gpl-3.0.html>.
          """)

def save_frame_as_slide(frame, output_folder, prs):
    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)    
    image_pil = Image.fromarray(frame_rgb)    
    image_stream = io.BytesIO()
    image_pil.save(image_stream, format='JPEG')
    image_stream.seek(0)    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)    
    left = top = Inches(1)
    slide.shapes.add_picture(image_stream, left, top, width=prs.slide_width - Inches(2), height=prs.slide_height - Inches(2))

def frame_diff(prev_frame, curr_frame):
    diff = cv2.absdiff(prev_frame, curr_frame)
    _, thresh = cv2.threshold(diff, 30, 255, cv2.THRESH_BINARY)
    return np.sum(thresh)

def get_unique_filename(output_folder, base_filename="DetectedSlides", extension=".pptx"):
    counter = 1
    unique_filename = f"{base_filename}{extension}"
    while os.path.exists(os.path.join(output_folder, unique_filename)):
        unique_filename = f"{base_filename}_{counter}{extension}"
        counter += 1
    return unique_filename

def process_video(video_path, output_folder, threshold, prs):
    file_name = input("Enter a file name for the presentation (without extension) and press [Enter] to initiate: ").strip()
    cap = cv2.VideoCapture(video_path)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    prev_frame = None
    for frame_index in tqdm(range(total_frames), desc="Processing video"):
        ret, frame = cap.read()
        if not ret:
            break        
        gray_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)        
        if prev_frame is not None:
            diff = frame_diff(prev_frame, gray_frame)
            if diff / (frame.size / 3) > threshold:
                save_frame_as_slide(frame, args.output_folder, prs)        
        prev_frame = gray_frame

    if not file_name:
        file_name = "Detected_Slides"
    unique_pptx_filename = get_unique_filename(args.output_folder, file_name)
    prs_file_path = os.path.join(args.output_folder, unique_pptx_filename)
    
    prs.save(prs_file_path)
    print(f"Presentation saved to {prs_file_path}")

    cap.release()

def main():
    print_gpl_notice()
    while True:
        proceed = input("Either type 'show w', 'show c', or proceed with processing the video? (y/n): ").strip().lower()
        if proceed == 'show w':
            show_warranty()
        elif proceed == 'show c':
            show_conditions()
        elif proceed == 'y':
            process_video(args.video_path, args.output_folder, args.threshold, prs)
            break
        elif proceed == 'n':
            print("Operation canceled.")
            sys.exit()
        else:
            print("Input invalid. Please try again.")
            continue

        proceed_with_processing = input("Proceed with processing the video? (y/n): ").strip().lower()
        if proceed_with_processing == 'y':
            process_video(args.video_path, args.output_folder, args.threshold, prs)
            break
        elif proceed_with_processing == 'n':
            print("Operation canceled.")
            sys.exit()
        else:
            print("Input invalid. Please try again.")

if __name__ == "__main__":
    main()