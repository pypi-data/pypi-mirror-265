from pptx import Presentation
from docx import Document
import re


class HandbookMaker:
    def __init__(self, ppt_file_path):
        """
        Initializes an instance of MyClass.

        Args:
            ppt_file_path (str): The file path of the PowerPoint presentation.

        Attributes:
            ppt_file_path (str): The file path of the PowerPoint presentation.
            ppt (Presentation): The PowerPoint presentation object.
            doc (Document): The document object.
            ebook_title_set (bool): Flag indicating if the ebook title has been set.
            slide_titles (list): List of slide titles.
        """
        self.ppt_file_path = ppt_file_path
        self.ppt = Presentation(ppt_file_path)
        self.doc = Document()
        self.ebook_title_set = False
        self.slide_titles = []

    def clean_text(self, text, skip_text_from_slides=[]):
        """
        Remove control characters and other non-XML compatible characters.

        Args:
            text (str): The input text to be cleaned.

        Returns:
            str: The cleaned text.

        """
        text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F]', '', text)
        text = re.sub(r'[\n]{2,}', '\n', text)
        text = re.sub(r'[\t]{2,}', '\t', text)
        for txt in skip_text_from_slides:
            text = re.sub(str(txt), "", text)
        return text

    def add_content_to_doc(self, content, is_header=False):
        if is_header:
            self.doc.add_heading(content, level=1)
        else:
            self.doc.add_paragraph(content)

    def convert_ppt_to_handbook(self, 
                                skip_slides:list[str]=None,
                                content_layout_name:str="Title and Content",
                                output_file_path:str="output.docx",
                                skip_text_from_slides:list[str]=[]):
        """
        Converts a PowerPoint presentation to a handbook document.

        Args:
            skip_slides (list): A list of slide layout names to skip during conversion.
            content_layout_name (str): The name of the slide layout to use for slide content.
            output_file_path (str): The file path to save the converted handbook document.

        Returns:
            str: The file path of the converted handbook document.
        """
        slide = self.ppt.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Module" in shape.text_frame.text:
                    self.doc.add_heading(shape.text_frame.text, level=0)
                    self.ebook_title_set = True
                    break

        for slide_idx, slide in enumerate(self.ppt.slides):
            if slide.slide_layout.name in skip_slides:
                continue
            if slide.slide_layout.name == "Section Header":
                self.doc.add_page_break()
                section_title = slide.shapes[2].text_frame.text
                self.doc.add_heading(section_title, level=1)
            else:
                if slide.slide_layout.name == content_layout_name and slide.shapes.title:
                    slide_title = slide.shapes.title.text
                    if slide_title not in self.slide_titles:
                        self.slide_titles.append(slide_title)
                        self.doc.add_page_break()
                        self.doc.add_heading(self.clean_text(slide_title,skip_text_from_slides), level=2)
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text
                    cleaned_notes_text = self.clean_text(notes_text,skip_text_from_slides)
                    self.add_content_to_doc(f"Notes: {cleaned_notes_text}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            cleaned_text = self.clean_text(paragraph.text,skip_text_from_slides)
                            self.add_content_to_doc(cleaned_text)

        self.doc.save(output_file_path)
        return output_file_path
