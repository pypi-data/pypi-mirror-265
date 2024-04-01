from pptx import Presentation
import requests
import json
import pdfplumber
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PPTxFile:
    def __init__(self, file_path):
        self.file_path = file_path
        self.prs = Presentation(self.file_path)
        self.slides = self.prs.slides
    
    def get_slide_text(self, slide):
        """Returns the text in a slide
        
        Arguments:
            slide {pptx.slide} -- slide object
        Returns:
            str -- text in slide
        """
        text_list = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # 12 corresponds to a group shape
                text_list.append("\n".join([subshape.text.strip() for subshape in shape.shapes if hasattr(subshape,"text")]).strip())
            else:
                if hasattr(shape, "text"):
                    text_list.append(shape.text)
        text_list = [text for text in text_list if text.strip()!=""]
        slide_text = '\n'.join(text_list)
        return slide_text

    def get_slide_notes(self, slide):
        """Returns the notes in a slide
        
        Arguments:
            slide {pptx.slide} -- slide object
        Returns:
            str -- notes in slide
        """
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text
        return notes_text

    def set_slide_notes(self, slide, text):
        """Sets the notes in a slide
        
        Arguments:
            slide {pptx.slide} -- slide object
            text {str} -- text to be set as notes
        """
        slide.notes_slide.notes_text_frame.text = text

    def erase_slide_notes(self, slide):
        slide.notes_slide.notes_text_frame.clear()

    def set_skip_slides(self, skip_slides_index: list[int] = None,
                        skip_slides_layout: list[str] = None,
                        skip_slide_titles: list[str] = None):
        """
        Set the slides to skip.
        :param skip_slides_index: Index of the slides to skip.
        :param skip_slides_layout: Layout of the slides to skip.
        :param skip_slide_titles: Titles of the slides to skip.
        """
        if skip_slides_index is None:
            skip_slides_index = [1, -1]
        if skip_slides_layout is None:
            skip_slides_layout = ["Title Slide", "CoverPage", "Quote Slide", 
                                  "Agenda", "Section Header", "QuoteHead"]
        if skip_slide_titles is None:
            skip_slide_titles = []
        self.skip_slides_index = skip_slides_index
        self.skip_slides_layout = skip_slides_layout
        self.skip_slide_titles = skip_slide_titles

    def convert_to_pdf(self, pspd_auth_token, output_path):
        """Converts the pptx file to pdf using PSPD API
        You need to generate an auth token from PSPDFKit website
        and pass it as an argument to this function
        Arguments:
            pdpd_auth_token {str} -- PSPDFKit auth token
            output_path {str} -- output path for the pdf file
        Returns:
            str -- status of the conversion
        """
        instructions = json.dumps({'parts': [{'file': 'document'}]})
        url = 'https://api.pspdfkit.com/build'
        self.pdf_path = output_path
        headers = {'Authorization': 'Bearer ' + pspd_auth_token}
        files = {'document': open(self.file_path, 'rb')}
        response = requests.request('POST', url, headers=headers,
                                    files=files,
                                    data={'instructions': instructions},
                                    stream=True)    
        if response.ok:
            with open(self.pdf_path, 'wb') as fd:
                for chunk in response.iter_content(chunk_size=8096):
                    fd.write(chunk)
            return "PDF file saved at: " + self.pdf_path
        else:
            return "PDF conversion failed."
    
    def get_text_from_pdf(self,
                          pdf_path: str = None,
                          omit_text: list = None):
        """Returns the text in the pdf file
        """
        if pdf_path is None:
            if self.pdf_path is None:
                print("PDF file path not provided")
                return None
            else:
                pdf_path = self.pdf_path

        all_text = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    # Replace newlines with spaces
                    if omit_text is not None:
                        for omit in omit_text:
                            text = text.replace(omit, '')
                all_text.append(text)
        return all_text
